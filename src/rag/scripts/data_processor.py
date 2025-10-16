"""
Data processing module for RAG system
Handles loading and processing various document formats
"""

import os
from pathlib import Path
from typing import List
from llama_index.core import Document
import pypdf
import docx2txt
from pptx import Presentation
import openpyxl
import pandas as pd
from tqdm import tqdm
from config import Config

class DataProcessor:
    """Process various document formats for RAG system"""
    
    SUPPORTED_FORMATS = {
        '.pdf': 'process_pdf',
        '.txt': 'process_txt',
        '.docx': 'process_docx',
        '.pptx': 'process_pptx',
        '.xlsx': 'process_xlsx',
        '.csv': 'process_csv',
        '.md': 'process_md'
    }
    
    def __init__(self, data_dir: str = None):
        self.data_dir = Path(data_dir).resolve() if data_dir else Config.DATA_DIR.resolve()
    
    def process_md(self, file_path: Path) -> str:
        """Extract text from Markdown file"""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                return file.read()
        except Exception as e:
            print(f"⚠️ Error processing MD {file_path}: {e}")
            return ""
    
    def process_pdf(self, file_path: Path) -> str:
        """Extract text from PDF file"""
        text = ""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = pypdf.PdfReader(file)
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
        except Exception as e:
            print(f"Error processing PDF {file_path.name}: {e}")
        return text
    
    def process_txt(self, file_path: Path) -> str:
        """Extract text from TXT file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except Exception as e:
            print(f"Error processing TXT {file_path}: {e}")
            return ""
    
    def process_docx(self, file_path: Path) -> str:
        """Extract text from DOCX file"""
        try:
            return docx2txt.process(str(file_path))
        except Exception as e:
            print(f"Error processing DOCX {file_path}: {e}")
            return ""
    
    def process_pptx(self, file_path: Path) -> str:
        """Extract text from PPTX file"""
        text = ""
        try:
            prs = Presentation(str(file_path))
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        except Exception as e:
            print(f"Error processing PPTX {file_path}: {e}")
        return text
    
    def process_xlsx(self, file_path: Path) -> str:
        """Extract text from XLSX file"""
        try:
            df = pd.read_excel(file_path)
            return df.to_string()
        except Exception as e:
            print(f"Error processing XLSX {file_path}: {e}")
            return ""
    
    def process_csv(self, file_path: Path) -> str:
        """Extract text from CSV file"""
        try:
            df = pd.read_csv(file_path)
            return df.to_string()
        except Exception as e:
            print(f"Error processing CSV {file_path}: {e}")
            return ""
    
    def load_documents(self) -> List[Document]:
        """Load all documents from data directory"""
        documents = []
        files = list(self.data_dir.rglob('*'))
        
        print(f"\n📂 Processing documents from: {self.data_dir}")
        
        for file_path in tqdm(files, desc="Loading documents"):
            if file_path.is_file():
                ext = file_path.suffix.lower()
                
                if ext in self.SUPPORTED_FORMATS:
                    processor_method = getattr(self, self.SUPPORTED_FORMATS[ext])
                    text = processor_method(file_path)
                    
                    if text.strip():
                        doc = Document(
                            text=text,
                            metadata={
                                'filename': file_path.name,
                                'file_path': str(file_path),
                                'file_type': ext
                            }
                        )
                        documents.append(doc)
        
        print(f"✓ Loaded {len(documents)} documents successfully")
        return documents

if __name__ == "__main__":
    # Test data processor
    processor = DataProcessor()
    docs = processor.load_documents()
    print(f"\nTotal documents loaded: {len(docs)}")
    if docs:
        print(f"Sample document: {docs[0].text[:200]}...")
